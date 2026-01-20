#!/usr/bin/env python3
"""
Generate PowerPoint presentation for GeneratePlaywrightAPITests project
Creates a comprehensive 10-slide presentation including frontend analysis
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

# Get project root
project_root = Path(__file__).parent.parent

def create_presentation():
    """Create PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "GeneratePlaywrightAPITests"
    subtitle.text = "Automated API Test Generation from OpenAPI/Swagger Specifications"
    
    # Slide 2: Problem & Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Problem & Solution"
    tf = content.text_frame
    tf.text = "The Challenge:"
    p = tf.add_paragraph()
    p.text = "• Manual API test creation is time-consuming and error-prone"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Maintaining tests as APIs evolve requires constant updates"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Test data generation often requires manual effort"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "\nThe Solution:"
    p = tf.add_paragraph()
    p.text = "Automated test generation from OpenAPI/Swagger specifications"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Python-based tool that generates Playwright API tests"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Creates realistic test data from schema definitions"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Handles resource dependencies and test ordering automatically"
    p.level = 1
    
    # Slide 3: Key Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Key Features"
    tf = content.text_frame
    tf.text = "Core Capabilities:"
    p = tf.add_paragraph()
    p.text = "✓ Automatic Test Generation"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   Parses OpenAPI/Swagger specs and generates tests for all endpoints"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "\n✓ Smart Resource Management"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   Creates resources (POST) first, stores IDs for dependent tests"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "\n✓ Flexible Data Generation"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   Schema-based (free) or AI-powered (OpenAI/Anthropic)"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "\n✓ Production Ready"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   Generates TypeScript Playwright test files ready for CI/CD"
    p.level = 2
    
    # Slide 4: How It Works
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "How It Works"
    tf = content.text_frame
    tf.text = "Simple 3-Step Process:"
    p = tf.add_paragraph()
    p.text = "1. Configuration"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   Configure Swagger URL in config.json"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "   Choose free (schema-based) or AI-powered mode"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "\n2. Generation"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   Run: python scripts/generate_tests_main_script.py"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "   Tool parses spec, extracts endpoints, generates test code"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "\n3. Execution"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   Run: npm run test"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "   View reports: npx playwright show-report"
    p.level = 2
    
    # Slide 5: Architecture Flow
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Architecture"
    tf = content.text_frame
    tf.text = "System Flow:"
    p = tf.add_paragraph()
    p.text = "OpenAPI/Swagger Specification"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   ↓"
    p = tf.add_paragraph()
    p.text = "Python Test Generator"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   • Parse spec • Extract endpoints • Generate test code"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "   ↓"
    p = tf.add_paragraph()
    p.text = "Playwright Test Files (.spec.ts)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   • TypeScript test code • Resource ID management"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "   ↓"
    p = tf.add_paragraph()
    p.text = "Test Execution & Reports"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   • Playwright runner • HTML reports • Coverage analysis"
    p.level = 2
    
    # Slide 6: Frontend Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Web Interface"
    tf = content.text_frame
    tf.text = "Modern React-Based Frontend:"
    p = tf.add_paragraph()
    p.text = "✓ Dark AI-themed interface"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ 5-step workflow (Connect → Generate → View → Run → Reports)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Real-time progress indicators"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Test code viewer with TypeScript/Gherkin formats"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "\nTech Stack:"
    p = tf.add_paragraph()
    p.text = "• React 19 + TypeScript"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Vite (fast build tool)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Tailwind CSS (styling)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Axios (HTTP client)"
    p.level = 1
    
    # Slide 7: Frontend Status
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Frontend Status"
    tf = content.text_frame
    tf.text = "Feature Completeness:"
    p = tf.add_paragraph()
    p.text = "✅ API Loading - 90% (Real API calls)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ Test Generation - 80% (Real generation)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ Test Viewing - 90% (Code viewer)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "❌ Test Execution - 0% (Mocked, needs backend)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "❌ Reports - 30% (Mock data)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "⚠️  Settings - 40% (UI only, no save)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "\nOverall: 61% Complete"
    p = tf.add_paragraph()
    p.text = "\nReady for backend integration!"
    p.level = 1
    
    # Slide 8: Frontend Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Frontend Features"
    tf = content.text_frame
    tf.text = "What's Working:"
    p = tf.add_paragraph()
    p.text = "✓ Loads Swagger/OpenAPI specs from URL"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Generates test code from API endpoints"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Displays code in TypeScript & Gherkin formats"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Copy & download test files"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Beautiful dark-themed UI"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "\nWhat's Needed:"
    p = tf.add_paragraph()
    p.text = "• Backend API for test execution"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• WebSocket for real-time logs"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Settings persistence"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Real reports integration"
    p.level = 1
    
    # Slide 9: Benefits
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Benefits"
    tf = content.text_frame
    tf.text = "Why Use This Tool?"
    p = tf.add_paragraph()
    p.text = "For Developers:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "⚡ Save hours of manual test writing"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "🔄 Keep tests in sync with API changes automatically"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "🎯 Focus on business logic, not boilerplate"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "\nFor QA Teams:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "📊 Comprehensive test coverage"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "🔍 Consistent test structure"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "📈 Easy CI/CD integration"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "\nFor Organizations:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "💰 Reduced testing costs • 🚀 Faster release cycles • ✅ Higher quality"
    p.level = 2
    
    # Slide 10: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Summary"
    subtitle.text = "✓ Automated API test generation from OpenAPI specs\n✓ Modern React web interface (61% complete)\n✓ Flexible - Free schema-based or AI-powered\n✓ Production-ready Playwright test files\n✓ Smart resource management\n✓ Easy to integrate into existing workflows\n\nNext: Backend API integration for test execution"
    
    # Save presentation
    output_path = project_root / "GeneratePlaywrightAPITests_Presentation.pptx"
    prs.save(str(output_path))
    print(f"✅ PowerPoint presentation created: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("❌ Error: python-pptx library not installed")
        print("📦 Installing python-pptx...")
        import subprocess
        import sys
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        print("✅ python-pptx installed. Running again...")
        create_presentation()
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
        import traceback
        traceback.print_exc()
